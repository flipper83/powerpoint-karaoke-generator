import os
import random

from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Emu, Pt

SLD_LAYOUT_TITLE = 0
SLD_LAYOUT_IMAGE = 6
SLD_LAYOUT_CHART = 6


def random_values(size):
    values = []
    max_value = 100
    for index_slide in range(0, size):
        current_value = random.randint(0, max_value)
        values.append(current_value / 100)
        max_value -= current_value
        if max_value < 0:
            max_value = 0

    return values


class KaraokePresentation:
    presentation = None

    def __init__(self):
        self.presentation = Presentation()

    def init(self, title_value, subtitle_value):
        title_slide_layout = self.presentation.slide_layouts[SLD_LAYOUT_TITLE]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_value
        subtitle.text = subtitle_value

    def add_pie_chart(self, categories):
        chart_layout = self.presentation.slide_layouts[SLD_LAYOUT_CHART]
        slide = self.presentation.slides.add_slide(chart_layout)

        chart_data = ChartData()
        chart_data.categories = categories
        values = random_values(len(categories))
        chart_data.add_series('Serie', values)

        self.add_title(slide)

        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '#%'

    def add_title(self, slide):
        x, y, cx, cy = Inches(2), Inches(0.5), Inches(6), Inches(4.5)
        title = slide.shapes.add_textbox(x, y, cx, cy)
        title.tsext = ""
        title.text_frame.paragraphs[0].font.size = Pt(50)
        title.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        title.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    def add_bar_chat(self, categories):
        chart_layout = self.presentation.slide_layouts[SLD_LAYOUT_CHART]
        slide = self.presentation.slides.add_slide(chart_layout)

        self.add_title(slide)

        chart_data = ChartData()
        chart_data.categories = categories
        values = random_values(len(categories))
        chart_data.add_series('Serie', values)

        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

    def add_image_slide(self, image_downloaded):
        image_layout = self.presentation.slide_layouts[SLD_LAYOUT_IMAGE]
        slide = self.presentation.slides.add_slide(image_layout)

        picture = self.picture_fitted(slide, image_downloaded)

        pos_x = (self.presentation.slide_width.inches - Emu(picture.width).inches) / 2
        picture.left = Inches(pos_x).emu

    def save(self):
        self.presentation.save('test.pptx')

    def get_image_size(self, image_downloaded):
        im = Image.open(image_downloaded)
        return im.size

    def picture_fitted(self, slide, image_downloaded):
        top = Inches(0)
        left = Inches(0)
        height = Inches(7.5)
        width = Inches(10)
        source_width, source_height = self.get_image_size(image_downloaded)
        if source_width > source_height:
            return slide.shapes.add_picture(image_downloaded, left, top, height=height)
        else:
            return slide.shapes.add_picture(image_downloaded, left, top, width=width)
